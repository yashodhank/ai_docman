"""Text and metadata extraction from various file types."""
from __future__ import annotations

import json
import mimetypes
import re
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Any


def get_file_metadata(path: Path) -> dict[str, Any]:
    """Extract file metadata using native Python + optional exiftool."""
    stat = path.stat()
    meta = {
        "filename": path.name,
        "extension": path.suffix.lower(),
        "size_bytes": stat.st_size,
        "created": datetime.fromtimestamp(stat.st_birthtime).isoformat()
        if hasattr(stat, "st_birthtime") else None,
        "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        "mime_type": mimetypes.guess_type(path.name)[0],
    }

    # Try exiftool for richer metadata (if available)
    try:
        result = subprocess.run(
            ["exiftool", "-json", "-n", str(path)],
            capture_output=True, text=True, timeout=10
        )
        if result.returncode == 0:
            exif = json.loads(result.stdout)[0]
            meta["exif"] = {
                k: v for k, v in exif.items()
                if k not in ("SourceFile", "Directory", "FileName")
            }
    except (FileNotFoundError, subprocess.TimeoutExpired, json.JSONDecodeError):
        pass

    return meta


def extract_text_pdf(path: Path, max_pages: int = 10) -> str:
    """Extract text from PDF using pdfplumber or PyMuPDF."""
    text_parts = []

    # Try pdfplumber first
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages[:max_pages]):
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        if text_parts:
            return "\n\n".join(text_parts)
    except ImportError:
        pass
    except Exception:
        pass

    # Fallback to PyMuPDF
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        for i, page in enumerate(doc[:max_pages]):
            text = page.get_text()
            if text.strip():
                text_parts.append(text)
        doc.close()
        if text_parts:
            return "\n\n".join(text_parts)
    except ImportError:
        pass
    except Exception:
        pass

    return ""


def extract_text_image(path: Path) -> str:
    """Extract text from image using Surya or EasyOCR."""
    # Try Surya first (faster, more accurate)
    try:
        from surya.ocr import run_ocr
        from surya.model.detection.model import load_model as load_det_model
        from surya.model.recognition.model import load_model as load_rec_model
        from PIL import Image

        det_model = load_det_model()
        rec_model = load_rec_model()
        image = Image.open(path)
        results = run_ocr([image], [["en"]], det_model, rec_model)
        if results and results[0].text_lines:
            return "\n".join(line.text for line in results[0].text_lines)
    except ImportError:
        pass
    except Exception:
        pass

    # Fallback to EasyOCR
    try:
        import easyocr
        reader = easyocr.Reader(["en"], gpu=False, verbose=False)
        results = reader.readtext(str(path))
        if results:
            return "\n".join(text for _, text, _ in results)
    except ImportError:
        pass
    except Exception:
        pass

    return ""


def extract_text_office(path: Path) -> str:
    """Extract text from Office documents."""
    ext = path.suffix.lower()

    # DOCX
    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            pass

    # XLSX
    if ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            texts = []
            for sheet in wb.sheetnames[:3]:  # First 3 sheets
                ws = wb[sheet]
                for row in list(ws.iter_rows(max_row=50, values_only=True)):
                    row_text = " ".join(str(c) for c in row if c)
                    if row_text.strip():
                        texts.append(row_text)
            return "\n".join(texts[:100])  # Limit
        except ImportError:
            pass

    # PPTX
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides[:20]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            pass

    return ""


def extract_text(path: Path, max_chars: int = 8000) -> str:
    """Extract text from any supported file type."""
    ext = path.suffix.lower()

    if ext == ".pdf":
        text = extract_text_pdf(path)
    elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"):
        text = extract_text_image(path)
    elif ext in (".docx", ".xlsx", ".pptx"):
        text = extract_text_office(path)
    elif ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".log"):
        try:
            text = path.read_text(errors="ignore")[:max_chars]
        except Exception:
            text = ""
    else:
        text = ""

    # Truncate and clean
    text = text[:max_chars]
    # Remove null bytes and other control characters that can break LLM parsing
    text = text.replace("\x00", " ")
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def analyze_file(path: Path) -> dict[str, Any]:
    """Full analysis of a file: metadata + text content."""
    meta = get_file_metadata(path)
    meta["text_preview"] = extract_text(path, max_chars=4000)
    meta["text_length"] = len(meta["text_preview"])
    return meta
